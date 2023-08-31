from datetime import date
from pathlib import Path
import json
from typing import Any, Dict

from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.oxml.shapes.groupshape import CT_GroupShape
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide

BUCKETS = [
    "WESTERN_NA",
    "CENTRAL_NA",
    "EASTERN_NA",
    "EUROPE",
    "INDIA",
    "AUSTRALIA_WEST",
    "AUSTRALIA_EAST",
    "NEW_ZEALAND",
    "INDONESIA",
    "PHILIPPINES",
    "SINGAPORE_MALAYSIA",
    "JAPAN_SOUTH_KOREA",
    "ISRAEL",
    "UAE",
    "BRAZIL",
    "COSTA_RICA",
]


def describe_slide(slide: Slide) -> None:
    for idx, shape in enumerate(slide.shapes):
        print(idx, shape)
        if hasattr(shape, "text"):
            print("\t", shape.text)

        if hasattr(shape, "shapes"):
            for idx2, shape2 in enumerate(shape.shapes):
                print("\t\t", idx2, shape2)
                if hasattr(shape2, "text"):
                    print("\t\t\t", shape2.text)


def edit_text(shape: Shape, new_text: Any) -> None:
    # If the text is 0, remove the parent shape
    if str(new_text) == "0":
        parent = shape._element.getparent()
        grandparent = parent.getparent()

        if isinstance(grandparent, GroupShape) or isinstance(
            grandparent, CT_GroupShape
        ):
            grandparent.remove(parent)
        else:
            parent.remove(shape._element)
        return

    runs = shape.text_frame.paragraphs[0].runs

    for run in runs:
        run.text = ""
    to_add = str(new_text)
    if len(to_add) == 1:
        to_add = " " + to_add
    runs[0].text = str(to_add)


TOLERANCE = 100_000


def move_shape(name: str, shape: GroupShape) -> Shape:
    location_data = json.loads(Path("locations.json").read_text())

    top = None
    left = None
    nested = None
    for location in location_data:
        if location["name"] == name:
            top = location["top"]
            left = location["left"]
            nested = location["nested"]
            break

    if top is None or left is None or nested is None:
        raise ValueError(f"Could not find location {name}")

    shape.top = top  # noqa
    shape.left = left  # noqa

    return list(shape.shapes)[1]


def get_shape_by_text(text: str, shapes: list[Shape]) -> Shape:
    for shape in shapes:
        if isinstance(shape, GroupShape):
            continue
        if shape.text_frame.paragraphs[0].runs[0].text == text:
            return shape

    raise ValueError(f"Could not find shape for {text}")


def populate_slide(buckets: Dict[str, int], date: date):
    ppt = Presentation("template-2023-08-31.pptx")
    # NOTE: ppt must have shape of 13.33 x 7.5 inches for the top and left to work

    slide = ppt.slides[0]

    shapes = list(slide.shapes)

    grouped_shapes = [shape for shape in shapes if isinstance(shape, GroupShape)]

    DATE = get_shape_by_text("AUGUST 24, 2023", shapes)
    date_str = date.strftime("%B %d, %Y")
    edit_text(DATE, date_str)i

    for shape, bucket in zip(grouped_shapes, BUCKETS):
        shape = move_shape(bucket, shape)
        edit_text(shape, buckets.pop(bucket))

    assert not buckets, buckets

    ppt.save("output.pptx")
